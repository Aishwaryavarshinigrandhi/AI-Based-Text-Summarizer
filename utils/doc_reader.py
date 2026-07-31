import os

import docx

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None


def extract_text_from_docx(file_path: str) -> str:
    if not os.path.exists(file_path):
        raise FileNotFoundError("DOCX file was not found")

    document = docx.Document(file_path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(paragraphs)


def extract_text_from_pptx(file_path: str) -> str:
    if not os.path.exists(file_path):
        raise FileNotFoundError("PPTX file was not found")

    if Presentation is None:
        raise RuntimeError("python-pptx is not installed")

    presentation = Presentation(file_path)
    slides = []
    for slide in presentation.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        if parts:
            slides.append("\n".join(part for part in parts if part))
    return "\n".join(slides)
